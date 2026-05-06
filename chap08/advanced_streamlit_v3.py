import streamlit as st
import os
import pandas as pd
import zipfile
from dotenv import load_dotenv
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import TextLoader, PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.runnables import RunnablePassthrough, RunnableParallel
from langchain_core.output_parsers import StrOutputParser
from langchain_community.chat_message_histories import StreamlitChatMessageHistory
from langchain_core.runnables.history import RunnableWithMessageHistory
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Pt
from io import BytesIO

# 1. 초기 설정
load_dotenv()
st.set_page_config(page_title="ITS-Guard Final Pro", page_icon="🛡️", layout="wide")
st.title("🛡️ ITS-Guard Final: 프리미엄 AI 어시스턴트")

if "model" not in st.session_state:
    st.session_state.model = ChatOpenAI(model="gpt-4o-mini", temperature=0)
msgs = StreamlitChatMessageHistory(key="final_pro_history")

# --- [고도화된 완성본 PPT 생성 함수] ---
def create_final_pptx(content_text):
    prs = Presentation()
    brand_color = RGBColor(0, 80, 136) # ITS-Guard Deep Blue
    accent_color = RGBColor(222, 255, 154) # 연두색 포인트

    # AI가 준 텍스트에서 섹션 추출
    sections = [s.strip() for s in content_text.split("###") if s.strip()]
    titles = [s.split(":")[0].strip() for s in sections if ":" in s]

    # --- [1. 제목 슬라이드] ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "ITS-Guard 프로젝트 결과 보고서"
    subtitle.text = "AI Road Safety Innovation\nSmart Infrastructure Management System"
    title.text_frame.paragraphs[0].font.color.rgb = brand_color
    title.text_frame.paragraphs[0].font.bold = True

    # --- [2. 목차 슬라이드] ---
    toc_slide = prs.slides.add_slide(prs.slide_layouts[1])
    toc_slide.shapes.title.text = "Table of Contents (목차)"
    tf = toc_slide.placeholders[1].text_frame
    tf.text = "본 제안서는 다음과 같은 순서로 구성되었습니다."
    for i, t in enumerate(titles):
        p = tf.add_paragraph()
        p.text = f"{i+1}. {t}"
        p.font.size = Pt(20)

    # --- [3. 본문 슬라이드] ---
    for section in sections:
        if ":" in section:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_part, body_part = section.split(":", 1)
            
            # 상단 디자인 바 추가 (수정된 부분!)
            header_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.5)
            )
            header_shape.fill.solid()
            # foreground_color -> fore_color 로 수정
            header_shape.fill.fore_color.rgb = brand_color 
            header_shape.line.fill.background()

            # 제목 스타일링
            slide.shapes.title.text = title_part.strip()
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            # 본문 내용
            tf = slide.placeholders[1].text_frame
            tf.text = body_part.strip()
            for para in tf.paragraphs:
                para.font.size = Pt(18)
                para.space_after = Pt(12)

    ppt_out = BytesIO()
    prs.save(ppt_out)
    ppt_out.seek(0)
    return ppt_out

# --- [RAG 엔진] ---
@st.cache_resource(show_spinner=False)
def get_retriever(uploaded_file):
    docs = TextLoader("./my_project.txt", encoding="utf-8").load()
    if uploaded_file is not None:
        with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_path = tmp_file.name
        loader = PyPDFLoader(tmp_path) if uploaded_file.name.endswith(".pdf") else TextLoader(tmp_path, encoding="utf-8")
        docs.extend(loader.load())
        os.remove(tmp_path)
    
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=100)
    split_docs = text_splitter.split_documents(docs)
    vectorstore = FAISS.from_documents(split_docs, OpenAIEmbeddings())
    return vectorstore.as_retriever(search_kwargs={"k": 3})

# --- [사이드바 도구] ---
with st.sidebar:
    st.header("📂 프리미엄 도구함")
    uploaded_file = st.file_uploader("학습용 문서 업로드", type=["pdf", "txt"])
    
    st.divider()
    if st.button("📝 대화 핵심 요약"):
        if len(msgs.messages) > 0:
            summary_chain = ChatPromptTemplate.from_messages([
                ("system", "대화 내용을 바탕으로 핵심 기술적 결론을 3줄 요약해줘."),
                MessagesPlaceholder(variable_name="history"),
                ("human", "요약해줘.")
            ]) | st.session_state.model | StrOutputParser()
            st.info(summary_chain.invoke({"history": msgs.messages}))
        else:
            st.warning("대화 내역이 없습니다.")

    if st.button("💎 완성본 PPT 제작"):
        if len(msgs.messages) > 0:
            with st.spinner("전문가용 디자인 슬라이드 구성 중..."):
                ppt_gen_chain = ChatPromptTemplate.from_messages([
                    ("system", """너는 세계 최고의 전략 컨설턴트야. 지금까지의 대화 내용을 바탕으로 
                    발표용 PPT 5~6장을 구성해줘.
                    
                    반드시 슬라이드마다 '### 제목: 내용' 형식을 지켜야 해.
                    내용은 불렛포인트(-)를 사용해서 가독성 좋게 적어줘.
                    
                    슬라이드 구성 예시:
                    1. 프로젝트 배경 및 필요성
                    2. 핵심 기술 엔진 (YOLOv8)
                    3. 실시간 분석 프로세스
                    4. 기대 효과 및 도입 이점
                    5. 향후 발전 로드맵"""),
                    MessagesPlaceholder(variable_name="history"),
                    ("human", "위 구성에 맞춰서 내용을 상세하게 짜줘.")
                ]) | st.session_state.model | StrOutputParser()
                
                ppt_content = ppt_gen_chain.invoke({"history": msgs.messages})
                ppt_file = create_final_pptx(ppt_content) # 업그레이드된 함수 호출
                
                st.download_button(
                    label="📥 디자인 PPT 다운로드",
                    data=ppt_file,
                    file_name="ITS_Guard_Premium_Report.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

    # 기존 '💎 완성본 PPT 제작' if문이 끝나는 지점 바로 아래에 추가하세요
    if st.button("📄 공식 기획안(Word) 제작"):
        if len(msgs.messages) > 0:
            with st.spinner("전문적인 기획안을 작성 중..."):
                # 기획안 전용 프롬프트
                report_gen_chain = ChatPromptTemplate.from_messages([
                    ("system", """대화 내용을 바탕으로 공식적인 프로젝트 기획안을 작성하세요. 
                    각 섹션은 반드시 '### 제목: 내용' 형식을 유지해야 합니다.
                    포함할 내용: 프로젝트 개요, 핵심 기술 스택, 시스템 아키텍처, 기대 효과."""),
                    MessagesPlaceholder(variable_name="history"),
                    ("human", "지금까지의 내용을 공식 기획안 형태로 정리해줘.")
                ]) | st.session_state.model | StrOutputParser()
                
                report_text = report_gen_chain.invoke({"history": msgs.messages})
                doc_file = create_report_docx(report_text) # 성용님이 작성하신 함수 호출
                
                st.download_button(
                    label="📥 기획안(Word) 다운로드",
                    data=doc_file,
                    file_name="ITS_Guard_Project_Plan.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )
        else:
            st.warning("대화 내역이 없어 기획안을 작성할 수 없습니다.")

    if st.button("📊 프로젝트 실행 예산/일정(Excel) 제작"):
        if len(msgs.messages) > 0:
            with st.spinner("데이터 시트 구성 중..."):
                excel_gen_chain = ChatPromptTemplate.from_messages([
                    ("system", """대화 내용을 바탕으로 프로젝트 실행에 필요한 리스트를 작성하세요.
                    형식은 반드시 '항목: 내용'으로 작성해야 엑셀에 깔끔하게 들어갑니다.
                    예: '서버 구축비: 500만원', '개발 기간: 3개월'"""),
                    MessagesPlaceholder(variable_name="history"),
                    ("human", "프로젝트 예산안과 일정을 엑셀용 리스트로 정리해줘.")
                ]) | st.session_state.model | StrOutputParser()
                
                excel_text = excel_gen_chain.invoke({"history": msgs.messages})
                excel_file = create_excel_file(excel_text)
                
                st.download_button(
                    label="📥 엑셀(Excel) 다운로드",
                    data=excel_file,
                    file_name="ITS_Guard_Execution_Plan.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )

    if st.button("🎁 프로젝트 통합 패키지(ZIP) 생성"):
        if len(msgs.messages) > 0:
            with st.spinner("모든 문서를 제작하고 압축하는 중..."):
                # 1. 각각의 내용을 생성하기 위한 AI 체인 정의 (빠른 생성을 위해 통합 프롬프트 사용)
                content_gen_chain = ChatPromptTemplate.from_messages([
                    ("system", "대화 내용을 바탕으로 PPT, Word, Excel에 들어갈 핵심 내용을 각각 '###' 구분자로 나누어 상세히 작성해줘."),
                    MessagesPlaceholder(variable_name="history"),
                    ("human", "통합 패키지 제작 시작.")
                ]) | st.session_state.model | StrOutputParser()
                
                full_content = content_gen_chain.invoke({"history": msgs.messages})
                
                # 2. 각 파일 객체 생성
                ppt_file = create_final_pptx(full_content)
                word_file = create_report_docx(full_content)
                excel_file = create_excel_file(full_content)
                
                # 3. ZIP 파일 생성
                zip_buffer = BytesIO()
                with zipfile.ZipFile(zip_buffer, "a", zipfile.ZIP_DEFLATED, False) as zip_file:
                    zip_file.writestr("01_발표자료_ITS_Guard.pptx", ppt_file.getvalue())
                    zip_file.writestr("02_상세계획서_ITS_Guard.docx", word_file.getvalue())
                    zip_file.writestr("03_실행예산및일정_ITS_Guard.xlsx", excel_file.getvalue())
                
                zip_buffer.seek(0)
                
                st.download_button(
                    label="📥 통합 패키지(.zip) 다운로드",
                    data=zip_buffer,
                    file_name="ITS_Guard_Full_Package.zip",
                    mime="application/zip"
                )
        else:
            st.warning("대화 내역이 없습니다.")

# --- [메인 채팅 엔진] ---
retriever = get_retriever(uploaded_file)
prompt = ChatPromptTemplate.from_messages([
    ("system", "너는 ITS-Guard 전문가야. 제공된 맥락을 바탕으로 전문적으로 답변해.\n\n맥락: {context}"),
    MessagesPlaceholder(variable_name="history"),
    ("human", "{question}"),
])

chain = (
    RunnableParallel({
        "context": (lambda x: x["question"]) | retriever,
        "question": lambda x: x["question"],
        "history": lambda x: x["history"]
    })
    .assign(answer=(RunnablePassthrough.assign(context=lambda x: x["context"]) | prompt | st.session_state.model | StrOutputParser()))
)

chain_with_history = RunnableWithMessageHistory(
    chain, lambda s: msgs, input_messages_key="question", history_messages_key="history", output_messages_key="answer"
)

# UI 출력
for msg in msgs.messages:
    st.chat_message(msg.type).write(msg.content)

if user_input := st.chat_input("프로젝트에 대해 질문하세요..."):
    st.chat_message("human").write(user_input)
    with st.spinner("AI 분석 중..."):
        res = chain_with_history.invoke({"question": user_input}, {"configurable": {"session_id": "pro_final"}})
        answer = res["answer"]
    with st.chat_message("ai"):
        st.write(answer)
        # TTS
        tts_url = f"https://dict.naver.com/api/tts?text={answer.replace(' ', '%20')[:200]}&speaker=mijin&speed=0"
        st.components.v1.html(f'<audio autoplay src="{tts_url}" />', height=0)
        with st.expander("📚 참고 문헌"):
            for doc in res["context"]:
                st.write(f"- {doc.metadata.get('source', '지식베이스')}: {doc.page_content[:150]}...")


def create_report_docx(report_content):
    doc = Document()
    # 보고서 제목
    title = doc.add_heading('ITS-Guard 프로젝트 상세 기획안', 0)
    
    # 본문 파싱 및 스타일 적용
    sections = report_content.split("###")
    for section in sections:
        if ":" in section:
            header, content = section.split(":", 1)
            doc.add_heading(header.strip(), level=1)
            p = doc.add_paragraph(content.strip())
            p.style.font.size = Pt(11)
        else:
            doc.add_paragraph(section.strip())
            
    doc_out = BytesIO()
    doc.save(doc_out)
    doc_out.seek(0)
    return doc_out

def create_excel_file(excel_content):
    # AI가 준 텍스트를 표 형태(리스트)로 파싱
    data = []
    lines = [l.strip() for l in excel_content.split('\n') if l.strip()]
    for line in lines:
        if ":" in line:
            data.append(line.split(":", 1))
    
    # 데이터프레임 생성 및 엑셀 변환
    df = pd.DataFrame(data, columns=["항목", "상세 내용"])
    excel_out = BytesIO()
    with pd.ExcelWriter(excel_out, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='Project_Plan')
    excel_out.seek(0)
    return excel_out