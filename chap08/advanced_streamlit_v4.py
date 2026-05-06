import streamlit as st
import os
import zipfile
import tempfile
import pandas as pd
from io import BytesIO
from dotenv import load_dotenv

# PPT, Word 관련 라이브러리
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Pt

# 랭체인 관련 라이브러리
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import TextLoader, PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.runnables import RunnablePassthrough, RunnableParallel
from langchain_core.output_parsers import StrOutputParser
from langchain_community.chat_message_histories import StreamlitChatMessageHistory
from langchain_core.runnables.history import RunnableWithMessageHistory
from langchain_core.tools import tool
from langchain_core.messages import ToolMessage

# 1. 초기 설정 및 환경 변수 로드
load_dotenv()
st.set_page_config(page_title="ITS-Guard Agent Pro", page_icon="🛡️", layout="wide")
st.title("🛡️ ITS-Guard Agent: 지무인형 AI 프로젝트 매니저")

if "model" not in st.session_state:
    st.session_state.model = ChatOpenAI(model="gpt-4o-mini", temperature=0)

msgs = StreamlitChatMessageHistory(key="agent_pro_history")

# --- [기존 문서 생성 함수들] ---

def create_final_pptx(content_text):
    prs = Presentation()
    brand_color = RGBColor(0, 80, 136)
    sections = [s.strip() for s in content_text.split("###") if s.strip()]
    
    # 제목 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "ITS-Guard 프로젝트 보고서"
    slide.placeholders[1].text = "AI 기반 스마트 도로 인프라 관리 시스템"
    
    for section in sections:
        if ":" in section:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title, body = section.split(":", 1)
            # 디자인 헤더 바
            header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.5))
            header.fill.solid()
            header.fill.fore_color.rgb = brand_color
            # 내용
            slide.shapes.title.text = title.strip()
            slide.placeholders[1].text = body.strip()
            
    ppt_out = BytesIO()
    prs.save(ppt_out)
    ppt_out.seek(0)
    return ppt_out

def create_report_docx(content_text):
    doc = Document()
    doc.add_heading('ITS-Guard 공식 프로젝트 기획안', 0)
    sections = content_text.split("###")
    for section in sections:
        if ":" in section:
            header, content = section.split(":", 1)
            doc.add_heading(header.strip(), level=1)
            doc.add_paragraph(content.strip())
    doc_out = BytesIO()
    doc.save(doc_out)
    doc_out.seek(0)
    return doc_out

def create_excel_file(content_text):
    data = []
    for line in content_text.split('\n'):
        if ":" in line:
            data.append(line.split(":", 1))
    df = pd.DataFrame(data, columns=["항목", "상세내용"])
    excel_out = BytesIO()
    with pd.ExcelWriter(excel_out, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='Plan')
    excel_out.seek(0)
    return excel_out

# --- [에이전트 도구(Tool) 정의] ---

@tool
def make_presentation(content: str):
    """대화 내용을 요약하여 전문가용 PPT 파일을 생성합니다."""
    return create_final_pptx(content)

@tool
def make_document(content: str):
    """대화 내용을 바탕으로 공식적인 Word 기획안을 생성합니다."""
    return create_report_docx(content)

@tool
def make_excel_sheet(content: str):
    """프로젝트 예산이나 일정표를 엑셀 파일로 생성합니다."""
    return create_excel_file(content)

tools = [make_presentation, make_document, make_excel_sheet]
llm_with_tools = st.session_state.model.bind_tools(tools)

# --- [RAG 엔진 설정] ---

@st.cache_resource(show_spinner=False)
def get_retriever(uploaded_file):
    docs = []
    if os.path.exists("./my_project.txt"):
        docs.extend(TextLoader("./my_project.txt", encoding="utf-8").load())
    
    if uploaded_file is not None:
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp.write(uploaded_file.getvalue())
            loader = PyPDFLoader(tmp.name) if uploaded_file.name.endswith(".pdf") else TextLoader(tmp.name, encoding="utf-8")
            docs.extend(loader.load())
    
    if not docs: return None
    
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=100)
    split_docs = text_splitter.split_documents(docs)
    vectorstore = FAISS.from_documents(split_docs, OpenAIEmbeddings())
    return vectorstore.as_retriever(search_kwargs={"k": 3})

# --- [UI 및 메인 로직] ---

with st.sidebar:
    st.header("📂 데이터 관리")
    uploaded_file = st.file_uploader("추가 지식 업로드", type=["pdf", "txt"])
    if st.button("🗑️ 대화 초기화"):
        msgs.clear()
        st.rerun()

# 채팅 출력
for msg in msgs.messages:
    if msg.type != "tool": # 도구 실행 메시지는 화면에 표시하지 않음
        st.chat_message(msg.type).write(msg.content)

if user_input := st.chat_input("질문하거나 'PPT 만들어줘'라고 명령하세요."):
    st.chat_message("human").write(user_input)
    
    with st.spinner("AI가 생각하고 행동하는 중..."):
        # 1. RAG 컨텍스트 가져오기
        retriever = get_retriever(uploaded_file)
        context = ""
        if retriever:
            docs = retriever.invoke(user_input)
            context = "\n".join([d.page_content for d in docs])

        # 2. 에이전트 판단 (메시지 이력 포함)
        system_prompt = f"너는 ITS-Guard 전문가야. 다음 맥락을 참고해: {context}\n파일 생성 요청 시 도구를 사용해."
        history = msgs.messages
        response = llm_with_tools.invoke([("system", system_prompt)] + history + [("human", user_input)])
        
        # 3. 도구 실행 여부 확인
        if response.tool_calls:
            msgs.add_user_message(user_input)
            msgs.add_ai_message(response)
            
            for tool_call in response.tool_calls:
                t_name = tool_call["name"]
                t_args = tool_call["args"]["content"]
                
                if t_name == "make_presentation":
                    file_data = make_presentation.invoke(tool_call["args"])
                    st.download_button("📥 PPT 다운로드", file_data, "ITS_Guard_Report.pptx")
                elif t_name == "make_document":
                    file_data = make_document.invoke(tool_call["args"])
                    st.download_button("📥 기획서 다운로드", file_data, "ITS_Guard_Plan.docx")
                elif t_name == "make_excel_sheet":
                    file_data = make_excel_sheet.invoke(tool_call["args"])
                    st.download_button("📥 예산일정표 다운로드", file_data, "ITS_Guard_Budget.xlsx")
                
                # 도구 실행 결과를 메시지에 추가 (AI에게 알림)
                msgs.add_message(ToolMessage(content="파일 생성 완료", tool_call_id=tool_call["id"]))
            
            st.success("요청하신 파일을 준비했습니다. 위 버튼을 클릭해 다운로드하세요!")
        else:
            # 1. 사용자 메시지 기록
            msgs.add_user_message(user_input)
            
            # 2. 스트리밍 답변 영역 생성
            with st.chat_message("ai"):
                # 08-4의 핵심: llm.stream()을 사용하여 한 글자씩 출력
                # 여기서는 도구가 바인딩되지 않은 기본 모델을 사용해 답변을 생성합니다.
                response_container = st.empty()
                full_response = ""
                
                # 스트리밍 루프
                for chunk in st.session_state.model.stream([("system", system_prompt)] + history + [("human", user_input)]):
                    full_response += chunk.content
                    response_container.markdown(full_response + "▌") # 타이핑 효과
                
                response_container.markdown(full_response) # 최종 결과 출력
                
            # 3. AI 답변 기록
            msgs.add_ai_message(full_response)