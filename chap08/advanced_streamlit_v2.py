import streamlit as st
import os
from dotenv import load_dotenv
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import TextLoader, PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.runnables import RunnablePassthrough, RunnableParallel
from langchain_core.output_parsers import StrOutputParser
from langchain_community.chat_message_histories import StreamlitChatMessageHistory
from langchain_core.runnables.history import RunnableWithMessageHistory
from pptx import Presentation
from io import BytesIO
import tempfile

# 1. 초기 설정
load_dotenv()
st.set_page_config(page_title="ITS-Guard Ultra", page_icon="🛡️")
st.title("🛡️ ITS-Guard Ultra: 멀티모달 가이드")

# 변수 사전 정의
model = ChatOpenAI(model="gpt-4o-mini", temperature=0)
msgs = StreamlitChatMessageHistory(key="ultra_chat")

# 2. RAG 엔진 함수 (업로드 파일 대응)
@st.cache_resource(show_spinner=False)
def get_retriever(uploaded_file):
    docs = TextLoader("./my_project.txt", encoding="utf-8").load()
    if uploaded_file is not None:
        with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_path = tmp_file.name
        loader = PyPDFLoader(tmp_path) if uploaded_file.name.endswith(".pdf") else TextLoader(tmp_path, encoding="utf-8")
        docs.extend(loader.load())
        os.remove(tmp_path)
    
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=100)
    split_docs = text_splitter.split_documents(docs)
    vectorstore = FAISS.from_documents(split_docs, OpenAIEmbeddings())
    return vectorstore.as_retriever(search_kwargs={"k": 3})

# 3. 사이드바 구성 (파일 업로드 및 요약)
with st.sidebar:
    st.header("📂 문서 학습 센터")
    uploaded_file = st.file_uploader("추가 기획서나 문서를 올려주세요", type=["pdf", "txt"])
    
    st.divider()
    if st.button("📝 대화 내용 요약하기"):
        if len(msgs.messages) > 0:
            summary_prompt = ChatPromptTemplate.from_messages([
                ("system", "너는 기술 회의 요약 전문가야. 핵심 기술적 쟁점과 결론을 3줄로 요약해줘."),
                MessagesPlaceholder(variable_name="history"),
                ("human", "지금까지의 대화 내용을 요약해줘.")
            ])
            summary_chain = summary_prompt | model | StrOutputParser()
            with st.status("요약 중..."):
                summary_result = summary_chain.invoke({"history": msgs.messages})
            st.info(summary_result)
        else:
            st.warning("요약할 내용이 없습니다.")

# 4. 메인 대화 로직
retriever = get_retriever(uploaded_file)
prompt = ChatPromptTemplate.from_messages([
    ("system", "너는 ITS-Guard 프로젝트의 AI 전문가야. 맥락을 바탕으로 친절하게 답변해줘.\n\n맥락: {context}"),
    MessagesPlaceholder(variable_name="history"),
    ("human", "{question}"),
])

chain = (
    RunnableParallel({
        "context": (lambda x: x["question"]) | retriever,
        "question": lambda x: x["question"],
        "history": lambda x: x["history"]
    })
    .assign(answer=(RunnablePassthrough.assign(context=lambda x: x["context"]) | prompt | model | StrOutputParser()))
)

chain_with_history = RunnableWithMessageHistory(
    chain, lambda s: msgs, input_messages_key="question", history_messages_key="history", output_messages_key="answer"
)

# 5. 채팅 출력 및 입력
for msg in msgs.messages:
    st.chat_message(msg.type).write(msg.content)

if user_input := st.chat_input("무엇이든 물어보세요..."):
    st.chat_message("human").write(user_input)
    with st.spinner("분석 중..."):
        res = chain_with_history.invoke({"question": user_input}, {"configurable": {"session_id": "ultra"}})
        answer = res["answer"]
        
    with st.chat_message("ai"):
        st.write(answer)
        # 🔊 자동 재생 TTS (네이버 TTS API 임시 사용)
        tts_html = f'<audio autoplay="true" src="https://dict.naver.com/api/tts?text={answer.replace(" ", "%20")[:200]}&speaker=mijin&speed=0" />'
        st.components.v1.html(tts_html, height=0)
        
        with st.expander("📚 참고 근거 확인"):
            for doc in res["context"]:
                st.write(f"- {doc.metadata.get('source', '기본지식')}: {doc.page_content[:150]}...")

# --- [새로 추가할 PPT 생성 함수] ---
def create_pptx(content_list):
    prs = Presentation()
    
    # 제목 슬라이드
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "ITS-Guard 프로젝트 제안서"
    slide.placeholders[1].text = "AI 자동 생성 가이드라인"

    # AI가 준 슬라이드별 내용을 바탕으로 페이지 생성
    for item in content_list:
        slide_layout = prs.slide_layouts[1] # 제목 + 본문 레이아웃
        slide = prs.slides.add_slide(slide_layout)
        
        # 텍스트가 "제목: 내용" 형태라고 가정하고 분리
        parts = item.split(":", 1)
        slide.shapes.title.text = parts[0].strip() if len(parts) > 0 else "핵심 내용"
        slide.placeholders[1].text = parts[1].strip() if len(parts) > 1 else item
        
    ppt_out = BytesIO()
    prs.save(ppt_out)
    ppt_out.seek(0)
    return ppt_out

# --- [사이드바에 버튼 추가] ---
with st.sidebar:
    st.divider()
    if st.button("📊 대화 기반 PPT 생성"):
        if len(msgs.messages) > 0:
            # AI에게 현재 대화 내용을 PPT 슬라이드 5개 분량으로 정리하라고 요청
            ppt_prompt = "지금까지의 대화 내용을 바탕으로 발표용 슬라이드 5장의 [제목: 내용] 구성을 짜줘."
            ppt_content = model.invoke(str(msgs.messages) + ppt_prompt).content
            
            # 내용을 리스트로 변환 (간단한 파싱)
            slides = ppt_content.split("\n")
            slides = [s for s in slides if ":" in s] # 형식이 맞는 것만 추출
            
            ppt_file = create_pptx(slides)
            
            st.download_button(
                label="📥 PPT 파일 다운로드",
                data=ppt_file,
                file_name="ITS_Guard_Presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        else:
            st.warning("대화 내용이 있어야 PPT를 만들 수 있습니다.")